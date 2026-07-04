#!/usr/bin/env python
"""
Generate CZT detector simulation tutorial PPTX.
Usage: python scripts/make_slides.py
Produces output/CZT_simulation_tutorial.pptx
"""
import os
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

W, H    = Inches(13.33), Inches(7.5)
MARGIN  = Inches(0.45)
TITLE_H = Inches(0.8)
BODY_TOP = MARGIN + TITLE_H + Inches(0.15)
BODY_W   = W - 2 * MARGIN
BODY_H   = H - BODY_TOP - MARGIN

BLUE   = RGBColor(0, 84, 166)
LBLUE  = RGBColor(220, 232, 248)
LGREY  = RGBColor(240, 240, 240)
DGREY  = RGBColor(110, 110, 110)
WHITE  = RGBColor(255, 255, 255)
BLACK  = RGBColor(0, 0, 0)
GREEN  = RGBColor(0, 130, 60)
LGREEN = RGBColor(220, 245, 228)
ORANGE = RGBColor(210, 95, 0)
LORAN  = RGBColor(255, 243, 220)

REPO = os.path.abspath(os.path.join(os.path.dirname(__file__), ".."))


# ── Helpers ───────────────────────────────────────────────────────────────────

def _prs():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    return prs

def _blank(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = WHITE
    return slide

def _title(slide, text):
    tb = slide.shapes.add_textbox(MARGIN, MARGIN, BODY_W, TITLE_H)
    tf = tb.text_frame
    p  = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = text
    run.font.bold = True
    run.font.size = Pt(26)
    run.font.color.rgb = BLUE
    bar = slide.shapes.add_shape(1, MARGIN, MARGIN + TITLE_H - Inches(0.05),
                                  BODY_W, Inches(0.05))
    bar.fill.solid(); bar.fill.fore_color.rgb = BLUE
    bar.line.fill.background()

def _tb(slide, text, left, top, width, height,
        size=13, bold=False, color=BLACK, align=PP_ALIGN.LEFT, wrap=True):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = wrap
    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color

def _code(slide, text, left, top, width, height, size=9):
    box = slide.shapes.add_shape(1, left, top, width, height)
    box.fill.solid(); box.fill.fore_color.rgb = LGREY
    box.line.color.rgb = DGREY; box.line.width = Pt(0.5)
    tf = box.text_frame; tf.word_wrap = False
    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        run = p.add_run()
        run.text = line
        run.font.name = "Courier New"
        run.font.size = Pt(size)
        run.font.color.rgb = BLACK

def _rect(slide, l, t, w, h, fill, line=None, lw=0.8):
    s = slide.shapes.add_shape(1, l, t, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = fill
    if line:
        s.line.color.rgb = line; s.line.width = Pt(lw)
    else:
        s.line.fill.background()
    return s

def _label(slide, text, l, t, w, h, size=10, bold=False,
           color=BLACK, align=PP_ALIGN.CENTER):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    p  = tf.paragraphs[0]; p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size); run.font.bold = bold
    run.font.color.rgb = color

def _img_placeholder(slide, label, l, t, w, h):
    _rect(slide, l, t, w, h, RGBColor(220, 230, 242), line=BLUE, lw=1.2)
    _label(slide, f"[ {label} ]", l, t + h // 2 - Inches(0.2),
           w, Inches(0.4), size=11, color=BLUE)

def _table(slide, headers, rows, left, top, col_widths, row_h=Inches(0.42)):
    col_x = [left]
    for cw in col_widths[:-1]:
        col_x.append(col_x[-1] + cw)
    for hdr, cx, cw in zip(headers, col_x, col_widths):
        _rect(slide, cx, top, cw, row_h, BLUE)
        _label(slide, hdr, cx + Inches(0.07), top + Inches(0.07),
               cw - Inches(0.14), row_h - Inches(0.1),
               size=10, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    for ri, row in enumerate(rows):
        bg = LGREY if ri % 2 == 0 else WHITE
        for cell, cx, cw in zip(row, col_x, col_widths):
            _rect(slide, cx, top + row_h * (ri + 1), cw, row_h, bg, line=DGREY, lw=0.3)
            _label(slide, cell, cx + Inches(0.07),
                   top + row_h * (ri + 1) + Inches(0.07),
                   cw - Inches(0.14), row_h - Inches(0.1),
                   size=10, color=BLACK, align=PP_ALIGN.LEFT)

def _bullet(slide, items, top, left=None, size=15, spacing=0.58):
    l = left if left is not None else MARGIN + Inches(0.2)
    t = top
    for item in items:
        _tb(slide, "•  " + item, l, t, BODY_W - Inches(0.4), Inches(0.5),
            size=size, color=BLACK)
        t += Inches(spacing)
    return t


# ── Slides ────────────────────────────────────────────────────────────────────

def slide_01_title(prs):
    slide = _blank(prs)
    _rect(slide, 0, 0, W, H, RGBColor(0, 50, 110))
    _tb(slide, "CZT Strip Detector Simulation",
        Inches(1.2), Inches(1.9), Inches(11.0), Inches(1.2),
        size=38, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    _tb(slide, "with SolidStateDetectors.jl",
        Inches(1.2), Inches(3.0), Inches(11.0), Inches(0.7),
        size=22, color=RGBColor(160, 200, 245), align=PP_ALIGN.CENTER)
    _tb(slide, "From geometry definition to electrode waveforms",
        Inches(1.2), Inches(3.9), Inches(11.0), Inches(0.5),
        size=15, color=RGBColor(140, 180, 230), align=PP_ALIGN.CENTER)


def slide_02_what_is_ssd(prs):
    slide = _blank(prs)
    _title(slide, "What is SolidStateDetectors.jl?")

    _bullet(slide, [
        "Julia package for simulating solid-state radiation detectors",
        "Solves the electric field and weighting potentials on an adaptive 3D grid",
        "Drifts charge carriers (e⁻ / h⁺) through the field, models trapping",
        "Computes induced signals on each electrode via the Shockley-Ramo theorem",
        "Geometry defined in a single YAML file — no mesh editor or C++ required",
    ], top=BODY_TOP + Inches(0.2))

    _rect(slide, MARGIN, BODY_TOP + Inches(3.5), BODY_W, Inches(1.1), LBLUE, line=BLUE)
    _tb(slide,
        "Shockley-Ramo theorem:   induced current on electrode k  =  q · v_drift · ∇Φ_w,k\n"
        "Φ_w,k  (weighting potential) is found by setting contact k to 1 V, all others to 0 V, "
        "then solving the Laplace equation.",
        MARGIN + Inches(0.2), BODY_TOP + Inches(3.62),
        BODY_W - Inches(0.4), Inches(0.9), size=13, color=BLUE)


def slide_03_two_files(prs):
    slide = _blank(prs)
    _title(slide, "Two Files, One Command")

    for i, (fname, desc, bg, lc) in enumerate([
        ("geometry.yaml",    "Defines the detector — crystal, electrodes, voltages, material",   LBLUE,  BLUE),
        ("single_event.jl",  "Solves fields, simulates charge drift, plots waveforms",            LGREEN, GREEN),
    ]):
        bx_t = BODY_TOP + Inches(0.3) + i * Inches(1.65)
        _rect(slide, MARGIN, bx_t, BODY_W, Inches(1.35), bg, line=lc)
        _tb(slide, fname, MARGIN + Inches(0.2), bx_t + Inches(0.12),
            Inches(4.5), Inches(0.55), size=20, bold=True, color=lc)
        _tb(slide, desc, MARGIN + Inches(0.2), bx_t + Inches(0.68),
            BODY_W - Inches(0.4), Inches(0.5), size=13, color=BLACK)

    _tb(slide, "Run with:", MARGIN, BODY_TOP + Inches(3.7),
        Inches(1.8), Inches(0.35), size=12, color=DGREY)
    _code(slide, "julia --project=. -t8 scripts/single_event.jl",
          MARGIN, BODY_TOP + Inches(4.1), BODY_W, Inches(0.55), size=12)
    _tb(slide, "-t8 sets the number of threads — more threads speeds up the field solve",
        MARGIN, BODY_TOP + Inches(4.8), BODY_W, Inches(0.4), size=11, color=DGREY)


def slide_04_geometry(prs):
    slide = _blank(prs)
    _title(slide, "The Geometry File")

    yaml = """\
semiconductor:
  material: CdZnTe          # pulls in built-in μ, ε tables
  temperature: 293          # K
  geometry:
    box:
      widths: [40, 40, 5]   # crystal: 40 × 40 × 5 mm  (z = drift axis)
  charge_drift_model:
    model: IsotropicChargeDriftModel
    mobilities: {e: 1000cm^2/(V*s), h: 50cm^2/(V*s)}
  charge_trapping_model:
    model: ConstantLifetime
    parameters: {τe: 1μs, τh: 1μs}   # good commercial CZT

contacts:
  - id: 1                   # integer label → indexes waveform output
    potential: 0            # applied voltage (V)
    geometry:
      box:
        hX: 0.05            # half-width → 100 µm full strip width
        hY: 20              # half-width → spans full 40 mm in y
        hZ: 0               # surface contact (no thickness)
        origin: [-2, 0, 2.5]  # centre of the electrode (mm)"""

    _code(slide, yaml, MARGIN, BODY_TOP, Inches(6.8), Inches(5.3), size=9)

    notes = [
        (Inches(0.0),  "material → built-in mobility & dielectric constants"),
        (Inches(0.65), "widths [x, y, z] — z is the 5 mm drift axis"),
        (Inches(1.3),  "IsotropicChargeDriftModel: constant μ (simplest option)"),
        (Inches(2.0),  "τe = τh = 1 µs → good commercial CZT,  (μτ)_e ~ 10⁻³ cm²/V"),
        (Inches(2.95), "each contact needs a unique id, potential, and geometry"),
        (Inches(3.6),  "hX/hY/hZ are half-widths — full width = 2× the value"),
        (Inches(4.25), "hZ: 0 → infinitesimally thin surface strip"),
        (Inches(4.9),  "origin is the centre point of the box, in mm"),
    ]
    ann_l = Inches(7.3)
    for dy, txt in notes:
        _rect(slide, ann_l - Inches(0.15), BODY_TOP + dy + Inches(0.12),
              Inches(0.08), Inches(0.08), BLUE)
        _tb(slide, txt, ann_l, BODY_TOP + dy,
            Inches(5.6), Inches(0.55), size=10, color=RGBColor(50, 50, 50))


def slide_05_refinement(prs):
    slide = _blank(prs)
    _title(slide, "The Adaptive Grid — Refinement Limits")

    # Left: concept flow boxes
    bx_w, bx_h = Inches(4.8), Inches(0.68)
    bx_l = MARGIN
    gap  = Inches(0.24)
    passes = [
        ("1  Initial coarse grid",                                       LGREY, BLACK),
        ("2  Refine where ΔΦ > 0.20 × V_bias → re-solve",              LBLUE, BLUE),
        ("3  Refine where ΔΦ > 0.10 × V_bias → denser near edges",     LBLUE, BLUE),
        ("4  Refine where ΔΦ > 0.05 × V_bias → finer at strips",       BLUE,  WHITE),
        ("5  Refine where ΔΦ > 0.02 × V_bias → ~14 µm at strip edge",  BLUE,  WHITE),
    ]
    y = BODY_TOP + Inches(0.1)
    for txt, bg, fg in passes:
        _rect(slide, bx_l, y, bx_w, bx_h, bg, line=BLUE, lw=0.8)
        _label(slide, txt, bx_l + Inches(0.15), y + Inches(0.12),
               bx_w - Inches(0.3), bx_h - Inches(0.2), size=11,
               bold=(bg == BLUE), color=fg, align=PP_ALIGN.LEFT)
        if txt != passes[-1][0]:
            _label(slide, "▼", bx_l + bx_w / 2 - Inches(0.1), y + bx_h,
                   Inches(0.25), gap + Inches(0.05), size=12, color=DGREY)
        y += bx_h + gap

    _rect(slide, bx_l, y + Inches(0.05), bx_w, Inches(0.52),
          RGBColor(235, 245, 228), line=GREEN, lw=0.8)
    _label(slide,
           "Grid adapts — fine only where field changes rapidly (near contacts)",
           bx_l + Inches(0.1), y + Inches(0.1), bx_w - Inches(0.2), Inches(0.38),
           size=10, color=GREEN, align=PP_ALIGN.LEFT)

    # Code snippet below flow
    _code(slide,
          "calculate_electric_potential!(sim;\n"
          "    refinement_limits  = [0.2, 0.1, 0.05, 0.02],\n"
          "    convergence_limit  = 1e-6,\n"
          "    depletion_handling = true)",
          bx_l, y + Inches(0.7), bx_w, Inches(1.0), size=9)

    # Right: real grid image
    grid_img = os.path.join(REPO, "output", "grid_refinement.png")
    img_l = Inches(5.5)
    img_w = W - img_l - MARGIN
    img_h = Inches(4.8)
    if os.path.exists(grid_img):
        slide.shapes.add_picture(grid_img, img_l, BODY_TOP, img_w, img_h)
    else:
        _img_placeholder(slide, "grid_refinement.png", img_l, BODY_TOP, img_w, img_h)

    # Spacing callout below image
    _rect(slide, img_l, BODY_TOP + img_h + Inches(0.08), img_w, Inches(0.52),
          LGREY, line=DGREY, lw=0.4)
    _label(slide,
           "At strip edge:  x-spacing 14 µm,  z-spacing 17 µm  (7 cells across 100 µm strip)\n"
           "In bulk:  x-spacing 1730 µm,  z-spacing 375 µm  — 35× coarser",
           img_l + Inches(0.1), BODY_TOP + img_h + Inches(0.1),
           img_w - Inches(0.2), Inches(0.45), size=9, color=DGREY, align=PP_ALIGN.LEFT)


def slide_06_charge_drift(prs):
    slide = _blank(prs)
    _title(slide, "Charge Cloud & Drift Knobs")

    _code(slide,
          "# Position must be in metres (not mm!)\n"
          "pos = CartesianPoint{Float32}(x_m, y_m, z_m)\n\n"
          "evt = Event([pos], [662.0u\"keV\"], N_CARRIERS;\n"
          "    radius           = [0.1u\"mm\"],   # physical cloud size\n"
          "    number_of_shells = 2,             # shell radii at 1r and 2r\n"
          ")\n\n"
          "simulate!(evt, sim;\n"
          "    Δt         = 0.1u\"ns\",    # waveform time step\n"
          "    max_nsteps = 50_000,       # max drift time = Δt × max_nsteps\n"
          ")\n\n"
          "wf = evt.waveforms[contact_id]   # .time and .signal per electrode",
          MARGIN, BODY_TOP, Inches(6.5), Inches(4.5), size=10)

    _table(slide,
           ["Knob", "Effect"],
           [
               ("radius",            "Cloud size in mm — default 0.5 mm causes staircase artifacts on small-pitch strips; use 0.1–0.2 mm"),
               ("N_CARRIERS",        "Macro-charge sampling points (not electron count) — more = smoother, slower"),
               ("number_of_shells",  "Shells placed at r, 2r, ... expanding the spatial extent of the cloud"),
               ("Δt",                "Time resolution of the output waveform"),
               ("max_nsteps",        "Max drift time = Δt × max_nsteps (50k × 0.1 ns = 5 µs)"),
               ("position (metres)", "CartesianPoint takes SI units — divide mm values by 1000"),
           ],
           Inches(6.8), BODY_TOP,
           [Inches(2.3), Inches(6.05)], row_h=Inches(0.48))


def slide_07_viz_fields(prs):
    slide = _blank(prs)
    _title(slide, "Built-in Visualizations — Fields & Geometry")

    _code(slide,
          "using Plots\n\n"
          "# 3D detector geometry with all contacts\n"
          "plot(sim.detector)\n\n"
          "# Electric potential — 2D cross-section\n"
          "plot(sim.electric_potential, y = 0.0u\"mm\")   # xz slice at y = 0\n\n"
          "# Electric field lines\n"
          "plot_electric_fieldlines!(sim, sampling = 3u\"mm\")\n\n"
          "# Weighting potential for one contact\n"
          "plot(sim.weighting_potentials[contact_id], y = 0.0u\"mm\")",
          MARGIN, BODY_TOP, Inches(6.0), Inches(3.8), size=10)

    # Key insight box
    _rect(slide, MARGIN, BODY_TOP + Inches(4.0), Inches(6.0), Inches(1.1),
          LBLUE, line=BLUE, lw=0.8)
    _tb(slide,
        "Weighting potential shape reveals the readout principle:\n"
        "• Small anode strip → WP flat across most of drift, rises steeply near strip\n"
        "  → electron signal only induced in final fraction of drift (small-pixel effect)\n"
        "• Broad cathode → WP linear → signal encodes interaction depth",
        MARGIN + Inches(0.15), BODY_TOP + Inches(4.1),
        Inches(5.7), Inches(1.0), size=10, color=BLUE)

    # Two images (or placeholders) stacked vertically on the right
    ph_w = Inches(5.8)
    ph_h = Inches(4.9)
    ep_img = os.path.join(REPO, "output", "electric_potential_xz.png")
    wp_img = os.path.join(REPO, "output", "weighting_potential_contact.png")
    for path, label, t_off in [
        (ep_img, "electric_potential_xz.png", BODY_TOP),
        (wp_img, "weighting_potential_contact.png", BODY_TOP + ph_h * 0.5),
    ]:
        if os.path.exists(path):
            slide.shapes.add_picture(path, Inches(6.6), t_off, ph_w, ph_h * 0.48)
        else:
            _img_placeholder(slide, label, Inches(6.6), t_off, ph_w, ph_h * 0.48)


def slide_08_viz_waveforms(prs):
    slide = _blank(prs)
    _title(slide, "Built-in Visualizations — Drift Paths & Waveforms")

    _code(slide,
          "# Overlay carrier drift trajectories on the geometry plot\n"
          "plot(sim.detector)\n"
          "plot!(evt.drift_paths)\n\n"
          "# Waveforms — all contacts, time vs induced charge\n"
          "plot(evt.waveforms)\n\n"
          "# Or access individually\n"
          "wf = evt.waveforms[contact_id]\n"
          "plot(wf.time, wf.signal)",
          MARGIN, BODY_TOP, Inches(5.8), Inches(3.0), size=10)

    _rect(slide, MARGIN, BODY_TOP + Inches(3.15), Inches(5.8), Inches(1.95),
          LGREY, line=DGREY, lw=0.5)
    rows = [
        ("Anode",   "Sharp rise as electrons collect → plateau height encodes energy",       BLUE),
        ("Cathode", "Slower, smaller amplitude → amplitude encodes interaction depth",        GREEN),
        ("Ratio",   "Cathode / anode amplitude → depth reconstruction without separate sensor", ORANGE),
    ]
    for ri, (label, desc, color) in enumerate(rows):
        y = BODY_TOP + Inches(3.25) + ri * Inches(0.58)
        _rect(slide, MARGIN + Inches(0.1), y, Inches(1.1), Inches(0.42), color)
        _label(slide, label, MARGIN + Inches(0.1), y, Inches(1.1), Inches(0.42),
               size=10, bold=True, color=WHITE)
        _tb(slide, desc, MARGIN + Inches(1.3), y + Inches(0.06),
            Inches(4.3), Inches(0.38), size=10, color=BLACK)

    drift_img = os.path.join(REPO, "output", "drift_paths.png")
    wf_img    = os.path.join(REPO, "output", "single_event.png")
    for path, label, t_off, h in [
        (drift_img, "drift_paths.png",           BODY_TOP,              Inches(2.45)),
        (wf_img,    "single_event.png waveforms", BODY_TOP + Inches(2.6), Inches(2.45)),
    ]:
        if os.path.exists(path):
            slide.shapes.add_picture(path, Inches(6.5), t_off, Inches(6.4), h)
        else:
            _img_placeholder(slide, label, Inches(6.5), t_off, Inches(6.4), h)


def slide_09_pitfalls(prs):
    slide = _blank(prs)
    _title(slide, "Common Pitfalls")

    pitfalls = [
        ("Geometry",      "hX/hY/hZ are half-widths — contacts end up 2× too large if you forget",                         BLUE),
        ("Geometry",      "hZ: 0 for surface contacts — any thickness changes the field solve",                             BLUE),
        ("Field solve",   "Solve order is fixed: E-potential → E-field → weighting potentials",                            ORANGE),
        ("Field solve",   "Solve time scales with number of contacts — cache the result if iterating on events",           ORANGE),
        ("Charge cloud",  "Default cloud radius 0.5 mm is too large for small-pitch strips → staircase artifact; use 0.1–0.2 mm", GREEN),
        ("Charge cloud",  "Position in CartesianPoint is in metres, not mm — divide your mm values by 1000",              GREEN),
        ("Material",      "τ_e = 10 µs (SSD.jl example) is research-grade CZT; good commercial = ~1 µs",                 RGBColor(120, 0, 120)),
        ("Output",        "Waveforms can be missing — always check before accessing .time / .signal",                      DGREY),
    ]

    row_h = Inches(0.54)
    t = BODY_TOP + Inches(0.05)
    for category, text, color in pitfalls:
        _rect(slide, MARGIN, t, Inches(1.4), row_h - Inches(0.06), color)
        _label(slide, category, MARGIN + Inches(0.05), t + Inches(0.06),
               Inches(1.3), row_h - Inches(0.18), size=9, bold=True, color=WHITE)
        _tb(slide, text, MARGIN + Inches(1.55), t + Inches(0.08),
            BODY_W - Inches(1.65), row_h - Inches(0.16), size=11, color=BLACK)
        t += row_h


def slide_10_flowchart(prs):
    slide = _blank(prs)
    _title(slide, "Simulation Pipeline — Overview")

    stages = [
        {
            "num":   "1",
            "title": "Load Geometry",
            "pseudo": "read YAML\n→ build detector with all contacts",
            "knob":  "geometry.yaml",
            "color": RGBColor(0, 100, 180),
        },
        {
            "num":   "2",
            "title": "Solve Fields",
            "pseudo": "solve Poisson eq. on adaptive grid → electric potential\n"
                      "take gradient → E-field\n"
                      "for each electrode: solve Laplace → weighting potential",
            "knob":  "refinement_limits\nconvergence_limit\ndepletion_handling",
            "color": RGBColor(0, 120, 80),
        },
        {
            "num":   "3",
            "title": "Simulate Event",
            "pseudo": "place charge cloud at interaction point\n"
                      "drift e⁻ and h⁺ through E-field, step by step\n"
                      "at each step: add induced charge on each electrode (Shockley-Ramo)",
            "knob":  "radius  ·  N_CARRIERS\nΔt  ·  max_nsteps",
            "color": RGBColor(150, 60, 0),
        },
        {
            "num":   "4",
            "title": "Read Out Waveforms",
            "pseudo": "for each electrode: induced charge vs time\nanode plateau → energy\ncathode amplitude → depth",
            "knob":  "waveforms[contact_id]",
            "color": RGBColor(100, 0, 140),
        },
    ]

    bx_w = Inches(2.8)
    bx_h = Inches(2.2)
    gap  = Inches(0.52)
    total = len(stages) * bx_w + (len(stages) - 1) * gap
    start_x = (W - total) / 2
    bx_y = BODY_TOP + Inches(0.3)

    for i, s in enumerate(stages):
        x = start_x + i * (bx_w + gap)
        c = s["color"]
        lc = RGBColor(min(c[0]+40,255), min(c[1]+40,255), min(c[2]+40,255))

        # main box
        _rect(slide, x, bx_y, bx_w, bx_h, c, line=WHITE, lw=0)

        # number badge
        _rect(slide, x + Inches(0.1), bx_y + Inches(0.1),
              Inches(0.38), Inches(0.38), WHITE)
        _label(slide, s["num"], x + Inches(0.1), bx_y + Inches(0.1),
               Inches(0.38), Inches(0.38), size=13, bold=True, color=c)

        # title
        _label(slide, s["title"],
               x + Inches(0.55), bx_y + Inches(0.1),
               bx_w - Inches(0.65), Inches(0.42),
               size=13, bold=True, color=WHITE, align=PP_ALIGN.LEFT)

        # pseudocode body
        _tb(slide, s["pseudo"],
            x + Inches(0.12), bx_y + Inches(0.58),
            bx_w - Inches(0.24), Inches(1.1),
            size=10, color=WHITE, wrap=True)

        # arrow between boxes
        if i < len(stages) - 1:
            ax = x + bx_w + Inches(0.08)
            ay = bx_y + bx_h / 2 - Inches(0.15)
            _rect(slide, ax, ay + Inches(0.1), gap - Inches(0.16),
                  Inches(0.04), DGREY)
            _label(slide, "▶", ax + gap - Inches(0.36), ay,
                   Inches(0.3), Inches(0.35), size=13, color=DGREY)

        # knob label below box
        _tb(slide, s["knob"],
            x, bx_y + bx_h + Inches(0.15),
            bx_w, Inches(0.55),
            size=9, color=DGREY, align=PP_ALIGN.LEFT if False else PP_ALIGN.CENTER)

    # slow/fast annotation
    _rect(slide, MARGIN, bx_y + bx_h + Inches(0.9), W - 2*MARGIN, Inches(0.55),
          LGREY, line=DGREY, lw=0.4)
    _label(slide,
           "Stage 2 (field solve) is slow — ~8–260 s depending on contacts and refinement.   "
           "Cache the result and reuse for all events on the same geometry.",
           MARGIN + Inches(0.15), bx_y + bx_h + Inches(0.95),
           W - 2*MARGIN - Inches(0.3), Inches(0.42),
           size=11, color=DGREY, align=PP_ALIGN.LEFT)


def slide_11_pseudocode(prs):
    slide = _blank(prs)
    _title(slide, "Simulation Pipeline — Code & Output")

    col_w = Inches(5.8)
    col2  = MARGIN + col_w + Inches(0.35)
    ph_w  = W - col2 - MARGIN

    # ── Left: pseudocode per stage ────────────────────────────────────────────
    stages = [
        ("1  Load geometry",
         "sim  ←  read_yaml(geometry.yaml)",
         LGREY, BLUE),
        ("2a  Electric potential",
         "for each grid point:\n"
         "  iterate until Δpotential < convergence_limit\n"
         "  → V(x,y,z) everywhere in crystal",
         LBLUE, BLUE),
        ("2b  Weighting potentials",
         "for each electrode k:\n"
         "  set electrode k = 1V, all others = 0V\n"
         "  solve Laplace  → Φ_w,k(x,y,z)",
         LBLUE, BLUE),
        ("3  Drift & induce",
         "place cloud at (x,y,z) with energy E\n"
         "repeat until carriers reach electrode:\n"
         "  move carriers by v = μ·E · Δt\n"
         "  induced_charge[k] += q · ΔΦ_w,k",
         RGBColor(255, 243, 220), ORANGE),
        ("4  Waveform output",
         "waveform[k]  =  induced_charge[k]  vs  time",
         RGBColor(240, 230, 248), RGBColor(100, 0, 140)),
    ]

    t = BODY_TOP + Inches(0.05)
    heights = [Inches(0.55), Inches(0.95), Inches(0.95), Inches(1.15), Inches(0.55)]
    for (label, pseudo, bg, fc), h in zip(stages, heights):
        _rect(slide, MARGIN, t, col_w, h, bg, line=fc, lw=1.0)
        _label(slide, label, MARGIN + Inches(0.12), t + Inches(0.06),
               col_w - Inches(0.24), Inches(0.32),
               size=10, bold=True, color=fc, align=PP_ALIGN.LEFT)
        if "\n" in pseudo or h > Inches(0.7):
            _tb(slide, pseudo,
                MARGIN + Inches(0.22), t + Inches(0.34),
                col_w - Inches(0.34), h - Inches(0.38),
                size=10, color=BLACK, wrap=True)
        else:
            _tb(slide, pseudo,
                MARGIN + Inches(0.22), t + Inches(0.16),
                col_w - Inches(0.34), h - Inches(0.2),
                size=10, color=BLACK)
        t += h + Inches(0.1)

    # ── Right: output placeholders ────────────────────────────────────────────
    ep_img  = os.path.join(REPO, "output", "electric_potential_xz.png")
    wp_img  = os.path.join(REPO, "output", "weighting_potentials.png")
    wf_img  = os.path.join(REPO, "output", "single_event.png")

    ph_h = Inches(1.95)
    gap  = Inches(0.15)
    img_t = BODY_TOP + Inches(0.05)

    for path, label in [(ep_img,  "electric_potential_xz.png"),
                        (wp_img,  "weighting_potentials.png"),
                        (wf_img,  "single_event.png — waveforms")]:
        if os.path.exists(path):
            slide.shapes.add_picture(path, col2, img_t, ph_w, ph_h)
        else:
            _img_placeholder(slide, label, col2, img_t, ph_w, ph_h)
        img_t += ph_h + gap


# ── Main ──────────────────────────────────────────────────────────────────────

def main():
    out = os.path.join(REPO, "output", "CZT_simulation_tutorial.pptx")
    os.makedirs(os.path.dirname(out), exist_ok=True)
    prs = _prs()
    slide_01_title(prs)
    slide_02_what_is_ssd(prs)
    slide_03_two_files(prs)
    slide_04_geometry(prs)
    slide_05_refinement(prs)
    slide_06_charge_drift(prs)
    slide_07_viz_fields(prs)
    slide_08_viz_waveforms(prs)
    slide_09_pitfalls(prs)
    slide_10_flowchart(prs)
    slide_11_pseudocode(prs)
    prs.save(out)
    print(f"Saved {len(prs.slides)} slides → {out}")

if __name__ == "__main__":
    main()
